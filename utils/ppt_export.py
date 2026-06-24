import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY  = RGBColor(0x1a, 0x1a, 0x2e)
BLUE  = RGBColor(0x6C, 0x8E, 0xF5)
GREEN = RGBColor(0x22, 0xc5, 0x5e)
RED   = RGBColor(0xf4, 0x3f, 0x5e)
GRAY  = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

STATUS_COLOR = {"done": GREEN, "ing": BLUE, "issue": RED, "plan": GRAY}
STATUS_LABEL = {
    "done":  "✅ 완료한 업무",
    "ing":   "🔄 진행중인 업무",
    "issue": "⚠️ 이슈 / 리스크",
    "plan":  "📅 다음주 계획",
}


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def _add_title_slide(prs, week_label, author):
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    tb = _textbox(slide, 1, 2.5, 8, 1.5)
    p = tb.text_frame.paragraphs[0]
    p.text = "주간 업무 보고서"
    p.font.size  = Pt(36)
    p.font.bold  = True
    p.font.color.rgb = WHITE
    p.alignment  = PP_ALIGN.CENTER

    tb2 = _textbox(slide, 1, 4.0, 8, 0.6)
    p2  = tb2.text_frame.paragraphs[0]
    p2.text = f"{week_label}   |   {author}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = BLUE
    p2.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs, key, items):
    slide = _blank_slide(prs)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = STATUS_COLOR[key]
    bar.line.fill.background()

    tb = _textbox(slide, 0.4, 0.15, 9, 0.8)
    p  = tb.text_frame.paragraphs[0]
    p.text = STATUS_LABEL[key]
    p.font.size  = Pt(22)
    p.font.bold  = True
    p.font.color.rgb = NAVY

    tb2 = _textbox(slide, 0.5, 1.1, 9, 5.5)
    tb2.text_frame.word_wrap = True
    for i, item in enumerate(items):
        para = tb2.text_frame.add_paragraph() if i > 0 else tb2.text_frame.paragraphs[0]
        para.text = f"• {item['title']}"
        para.font.size  = Pt(14)
        para.font.bold  = True
        para.font.color.rgb = NAVY
        para.space_after = Pt(2)
        if item.get("description"):
            sub = tb2.text_frame.add_paragraph()
            sub.text = f"   └ {item['description']}"
            sub.font.size = Pt(11)
            sub.font.color.rgb = RGBColor(0x71, 0x80, 0x96)
            sub.space_after = Pt(10)


def _add_summary_slide(prs, sections):
    slide = _blank_slide(prs)

    positions = [
        (0.2, 0.1), (5.2, 0.1),
        (0.2, 3.75), (5.2, 3.75),
    ]
    for (x, y), key in zip(positions, ["done", "ing", "issue", "plan"]):
        tb = _textbox(slide, x, y, 4.6, 3.4)
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = STATUS_LABEL[key]
        p.font.size  = Pt(13)
        p.font.bold  = True
        p.font.color.rgb = STATUS_COLOR[key]
        p.space_after = Pt(6)

        for item in sections[key]:
            para = tf.add_paragraph()
            para.text = f"• {item['title']}"
            para.font.size = Pt(11)
            para.font.color.rgb = NAVY

        if not sections[key]:
            para = tf.add_paragraph()
            para.text = "해당 없음"
            para.font.size = Pt(11)
            para.font.color.rgb = GRAY


def create_ppt(work_list: list, week_label: str, author: str = "", layout: str = "summary") -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    sections = {"done": [], "ing": [], "issue": [], "plan": []}
    for w in work_list:
        sections[w["status"]].append(w)

    _add_title_slide(prs, week_label, author)

    if layout == "summary":
        _add_summary_slide(prs, sections)
    else:
        for key in ["done", "ing", "issue", "plan"]:
            if sections[key]:
                _add_section_slide(prs, key, sections[key])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()